from __future__ import annotations

import ast
import json
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import pytest
from leo_ppt_generator.contracts import ContractError, PageArtifact
from leo_ppt_generator.editable.adapter import EditableAdapter
from leo_ppt_generator.image_deck.adapter import ImageDeckAdapter
from PIL import Image
from pptx import Presentation

from tests.ppt_fixtures import build_editable_page


def _png(path: Path, color: str = "navy") -> Path:
    Image.new("RGB", (160, 90), color).save(path)
    return path


def test_page_artifact_version_and_source_hash_fail_closed(tmp_path):
    source = _png(tmp_path / "source.png")
    artifact = PageArtifact.from_source("page_001", "image", source, source, None)
    artifact.verify()
    with pytest.raises(ContractError, match="unknown_contract_version"):
        PageArtifact(**{**artifact.to_dict(), "schema_version": 99})
    source.write_bytes(b"changed")
    with pytest.raises(ContractError, match="source_hash_mismatch"):
        artifact.verify()


def test_editable_page_artifact_requires_validation_reference(tmp_path):
    source = _png(tmp_path / "source.png")
    page, manifest = build_editable_page(source, tmp_path / "page.pptx")
    artifact = PageArtifact.from_source(
        "page_001", "editable", source, page, None, manifest=manifest
    )

    with pytest.raises(ContractError, match="validation_missing"):
        artifact.verify()


def test_image_adapter_rejects_missing_slide_and_builds_complete_fixture(tmp_path):
    deck = tmp_path / "deck"
    adapter = ImageDeckAdapter(deck)
    adapter.prepare([{"number": 1, "notes": "one"}, {"number": 2, "notes": "two"}])
    adapter.record(1, _png(tmp_path / "one.png", "red"), backend="fixture", expected_revision=0, operation_id="record-1")
    with pytest.raises(ContractError, match="missing_page_artifact"):
        adapter.finalize(tmp_path / "incomplete.pptx")
    adapter.record(2, _png(tmp_path / "two.png", "blue"), backend="fixture", expected_revision=1, operation_id="record-2")
    result = adapter.finalize(tmp_path / "complete.pptx")
    assert result["status"] == "completed"
    assert Path(result["pptx"]).is_file()
    Path(result["pptx"]).write_bytes(b"tampered")
    with pytest.raises(ContractError, match="image_assemble_rebuild_required"):
        adapter.finalize(tmp_path / "complete.pptx")


def test_image_adapter_serializes_concurrent_records_and_replays_same_operation(tmp_path):
    adapter = ImageDeckAdapter(tmp_path / "deck")
    adapter.prepare([{"number": 1}, {"number": 2}])
    sources = [_png(tmp_path / "one.png", "red"), _png(tmp_path / "two.png", "blue")]

    def record(number):
        try:
            return adapter.record(number, sources[number - 1], backend="fixture", expected_revision=0, operation_id=f"record-{number}")
        except ContractError as error:
            return str(error)

    with ThreadPoolExecutor(max_workers=2) as pool:
        results = list(pool.map(record, (1, 2)))
    assert sum(isinstance(result, PageArtifact) for result in results) == 1
    assert results.count("vendor_revision_conflict") == 1
    completed_number = 1 if isinstance(results[0], PageArtifact) else 2
    replay = adapter.record(completed_number, sources[completed_number - 1], backend="fixture", expected_revision=0, operation_id=f"record-{completed_number}")
    assert replay.page_id == f"page_{completed_number:03d}"


def test_domain_adapters_reject_mutations_after_run_cancel(tmp_path):
    run = tmp_path / "run"
    (run / "run.json").parent.mkdir(parents=True)
    (run / "run.json").write_text('{"status": "cancelled"}\n', encoding="utf-8")
    image = ImageDeckAdapter(run / "image-deck")
    image.prepare([{"number": 1}])
    with pytest.raises(ContractError, match="run_cancelled_mutation_forbidden"):
        image.record(1, _png(tmp_path / "cancelled.png"), backend="fixture", expected_revision=0, operation_id="cancelled")


def test_editable_adapter_blocks_multi_page_without_worker(tmp_path):
    adapter = EditableAdapter(tmp_path / "editable")
    result = adapter.prepare([_png(tmp_path / "a.png"), _png(tmp_path / "b.png")], worker_available=False)
    assert result["status"] == "blocked"
    assert result["reason_code"] == "worker_capability_unavailable"


def test_editable_prepare_persists_source_notes_for_finalize_contract(tmp_path):
    source = _png(tmp_path / "source.png")
    adapter = EditableAdapter(tmp_path / "editable")

    adapter.prepare([source], worker_available=False, notes={1: "speaker note"})

    assert adapter._jobs()["pages"][0]["notes"] == "speaker note"


def test_editable_dispatch_request_honors_configured_concurrency(tmp_path, monkeypatch):
    monkeypatch.setenv("LEO_PPT_MAX_WORKERS", "2")
    adapter = EditableAdapter(tmp_path / "editable")
    sources = [_png(tmp_path / f"page-{number}.png") for number in range(1, 4)]
    result = adapter.prepare(sources, worker_available=True)

    assert result["next_action"] == {
        "kind": "request_worker_dispatch",
        "payload": {
            "dispatch_requirement": "multi_agent_required",
            "page_count": 3,
            "estimated_duration_per_page_seconds": 180,
            "suggested_max_concurrent": 2,
            "runtime_fallback": False,
        },
    }


def test_editable_dispatch_rejects_worker_id_path_traversal(tmp_path):
    adapter = EditableAdapter(tmp_path / "editable")
    adapter.prepare([_png(tmp_path / "source.png")], worker_available=False)
    prompt = tmp_path / "prompt.md"
    prompt.write_text("worker", encoding="utf-8")

    with pytest.raises(ContractError, match="invalid_agent_id"):
        adapter.dispatch("page_001", "../../escape", prompt)

    assert not (tmp_path / "escape").exists()


def test_editable_adapter_records_and_reverifies_domain_artifacts(tmp_path):
    source = _png(tmp_path / "source.png")
    adapter = EditableAdapter(tmp_path / "editable")
    prepared = adapter.prepare([source], worker_available=False)
    assert prepared["reason_code"] == "single_unit_current_agent_allowed"
    page, manifest = build_editable_page(source, tmp_path / "page.pptx", text="Editable")
    validation = tmp_path / "validation.json"
    validation.write_text('{"passed": true}\n', encoding="utf-8")
    artifact = adapter.record("page_001", page, validation, manifest, expected_revision=0, operation_id="record-1", notes="speaker")
    assert artifact.mode == "editable"
    assert adapter.artifacts()[0].notes == "speaker"
    replay = adapter.record("page_001", page, validation, manifest, expected_revision=0, operation_id="record-1", notes="speaker")
    assert replay.artifact_sha256 == artifact.artifact_sha256
    validation.write_text('{"passed": true, "changed": true}\n', encoding="utf-8")
    with pytest.raises(ContractError, match="validation_hash_mismatch"):
        adapter.artifacts()
    with pytest.raises(ContractError, match="vendor_revision_conflict"):
        adapter.record("page_001", page, validation, manifest, expected_revision=0, operation_id="record-2")


def test_editable_adapter_rejects_stale_passed_validation(tmp_path):
    source = _png(tmp_path / "source.png")
    adapter = EditableAdapter(tmp_path / "editable")
    adapter.prepare([source], worker_available=False)
    page, manifest = build_editable_page(source, tmp_path / "page.pptx", text="Original")
    manifest_value = json.loads(manifest.read_text(encoding="utf-8"))
    manifest_value["text_boxes"][0]["text"] = "Changed after validation"
    manifest_value["text_inventory"] = ["Changed after validation"]
    manifest.write_text(json.dumps(manifest_value), encoding="utf-8")
    validation = tmp_path / "validation.json"
    validation.write_text('{"passed": true}\n', encoding="utf-8")

    with pytest.raises(ContractError, match="page_validation_failed"):
        adapter.record(
            "page_001",
            page,
            validation,
            manifest,
            expected_revision=0,
            operation_id="stale-validation",
        )


def test_editable_finalize_rejects_a_tampered_delivery_on_replay(tmp_path):
    source = _png(tmp_path / "source.png")
    adapter = EditableAdapter(tmp_path / "editable")
    adapter.prepare([source], worker_available=False)
    page, manifest = build_editable_page(source, tmp_path / "page.pptx")
    validation = tmp_path / "validation.json"
    validation.write_text('{"passed": true}\n', encoding="utf-8")
    adapter.record(
        "page_001",
        page,
        validation,
        manifest,
        expected_revision=0,
        operation_id="record-1",
    )
    first = adapter.finalize(tmp_path / "deck.pptx")
    Path(first["pptx"]).write_bytes(b"tampered")

    with pytest.raises(ContractError, match="editable_finalize_manifest_conflict"):
        adapter.finalize(tmp_path / "deck.pptx")


def test_editable_finalize_preserves_custom_source_canvas(tmp_path):
    source = tmp_path / "source.png"
    Image.new("RGB", (400, 300), "navy").save(source)
    adapter = EditableAdapter(tmp_path / "editable")
    adapter.prepare([source], worker_available=False)
    page, manifest = build_editable_page(
        source,
        tmp_path / "page.pptx",
        slide_width=400 / 96,
        slide_height=300 / 96,
    )
    validation = tmp_path / "validation.json"
    validation.write_text('{"passed": true}\n', encoding="utf-8")
    artifact = adapter.record(
        "page_001",
        page,
        validation,
        manifest,
        expected_revision=0,
        operation_id="record-custom-canvas",
    )
    result = adapter.finalize(tmp_path / "custom.pptx")
    presentation = Presentation(result["pptx"])

    assert (artifact.width, artifact.height) == (400, 300)
    assert presentation.slide_width / presentation.slide_height == pytest.approx(4 / 3)


def test_editable_adapter_rejects_failed_validation_and_non_single_page(tmp_path):
    source = _png(tmp_path / "source.png")
    adapter = EditableAdapter(tmp_path / "editable")
    adapter.prepare([source], worker_available=False)
    page, manifest = build_editable_page(source, tmp_path / "page.pptx")
    failed = tmp_path / "failed.json"
    failed.write_text('{"passed": false}\n', encoding="utf-8")
    with pytest.raises(ContractError, match="page_validation_failed"):
        adapter.record("page_001", page, failed, manifest, expected_revision=0, operation_id="failed-record")
    multi = Presentation(page)
    multi.slides.add_slide(multi.slide_layouts[0])
    multi.save(page)
    passed = tmp_path / "passed.json"
    passed.write_text('{"passed": true}\n', encoding="utf-8")
    with pytest.raises(ContractError, match="invalid_editable_page"):
        adapter.record("page_001", page, passed, manifest, expected_revision=0, operation_id="multi-record")


def test_non_vendor_modules_do_not_import_vendor_directly():
    package_root = Path(__file__).resolve().parents[2] / "skills/leo-ppt-generator/runtime/src/leo_ppt_generator"
    allowed = {package_root / "image_deck/adapter.py", package_root / "editable/adapter.py"}
    offenders = []
    for path in package_root.rglob("*.py"):
        if "_vendor" in path.parts or path in allowed:
            continue
        tree = ast.parse(path.read_text(encoding="utf-8"))
        for node in ast.walk(tree):
            if isinstance(node, ast.ImportFrom) and node.module and "_vendor" in node.module:
                offenders.append(path)
            if isinstance(node, ast.Import) and any("_vendor" in item.name for item in node.names):
                offenders.append(path)
    assert offenders == []
