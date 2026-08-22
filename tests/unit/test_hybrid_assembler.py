from __future__ import annotations

from pathlib import Path

import pytest
from leo_ppt_generator.contracts import ContractError, PageArtifact
from leo_ppt_generator.hybrid.assembler import HybridAssembler
from PIL import Image
from pptx import Presentation

from tests.ppt_fixtures import build_editable_page


def _image(path: Path, color: str) -> Path:
    Image.new("RGB", (320, 180), color).save(path)
    return path


def _artifact(tmp_path: Path, number: int, mode: str) -> PageArtifact:
    source = _image(tmp_path / f"source-{number}.png", f"#{number}{number}{number}{number}{number}{number}")
    if mode == "image":
        artifact, manifest = source, None
    else:
        artifact, manifest = build_editable_page(source, tmp_path / f"page-{number}.pptx", text=f"page {number}")
    validation = tmp_path / f"validation-{number}.json"
    validation.write_text('{"passed": true}\n', encoding="utf-8")
    return PageArtifact.from_source(f"page_{number:03d}", mode, source, artifact, validation, manifest=manifest, notes=f"notes {number}")


def test_hybrid_preserves_page_order_modes_and_notes(tmp_path):
    artifacts = [_artifact(tmp_path, 1, "editable"), _artifact(tmp_path, 2, "image"), _artifact(tmp_path, 3, "editable")]
    result = HybridAssembler().assemble(artifacts, tmp_path / "hybrid.pptx", selected_pages={1, 3})
    assert result["delivery_type"] == "hybrid"
    assert [page["mode"] for page in result["pages"]] == ["editable", "image", "editable"]
    assert len(Presentation(result["pptx"]).slides) == 3


def test_partial_hybrid_requires_confirmation_bound_to_current_failure_set(tmp_path):
    artifacts = [_artifact(tmp_path, 1, "editable"), _artifact(tmp_path, 2, "image")]
    failures = {2: "page_validation_failed"}
    with pytest.raises(ContractError, match="partial_hybrid_confirmation_required"):
        HybridAssembler().assemble(artifacts, tmp_path / "blocked.pptx", selected_pages={1, 2}, failures=failures)
    result = HybridAssembler().assemble(
        artifacts,
        tmp_path / "partial.pptx",
        selected_pages={1, 2},
        failures=failures,
        partial_confirmation=HybridAssembler.failure_fingerprint(
            failures,
            selected_pages={1, 2},
            baseline_fingerprint=HybridAssembler.baseline_fingerprint(artifacts),
        ),
    )
    assert result["delivery_type"] == "partial-hybrid"
    assert result["failures"] == {"2": "page_validation_failed"}


def test_same_hybrid_request_replays_without_overwriting_delivery(tmp_path):
    artifacts = [_artifact(tmp_path, 1, "editable")]
    output = tmp_path / "deck.pptx"
    first = HybridAssembler().assemble(artifacts, output)
    second = HybridAssembler().assemble(artifacts, output)
    assert first["idempotency_status"] == "created"
    assert second["idempotency_status"] == "replayed"
    assert second["pptx_sha256"] == first["pptx_sha256"]


def test_hybrid_preserves_matching_custom_canvas_and_rejects_mixed_canvas(tmp_path):
    custom_source = tmp_path / "custom-source.png"
    Image.new("RGB", (400, 300), "navy").save(custom_source)
    custom_page, custom_manifest = build_editable_page(
        custom_source,
        tmp_path / "custom-page.pptx",
        slide_width=400 / 96,
        slide_height=300 / 96,
    )
    validation = tmp_path / "custom-validation.json"
    validation.write_text('{"passed": true}\n', encoding="utf-8")
    custom = PageArtifact.from_source(
        "page_001",
        "editable",
        custom_source,
        custom_page,
        validation,
        manifest=custom_manifest,
    )

    result = HybridAssembler().assemble([custom], tmp_path / "custom-hybrid.pptx")
    presentation = Presentation(result["pptx"])
    assert presentation.slide_width / presentation.slide_height == pytest.approx(4 / 3)

    wide = _artifact(tmp_path, 2, "image")
    with pytest.raises(ContractError, match="page_size_mismatch"):
        HybridAssembler().assemble(
            [custom, wide],
            tmp_path / "mixed.pptx",
            selected_pages={1},
        )
