from __future__ import annotations

import asyncio
import json
import sys
from pathlib import Path
from types import SimpleNamespace
from urllib.request import Request

import fitz
from leo_ppt_generator._vendor.codex_ppt.assemble_ppt import create_presentation
from leo_ppt_generator._vendor.codex_ppt.image_providers.atlascloud import (
    AtlasCloudImageProvider,
)
from leo_ppt_generator._vendor.codex_ppt.image_providers.openai_compatible import (
    OpenAICompatibleImageProvider,
)
from leo_ppt_generator._vendor.codex_ppt.remove_chroma_key import (
    _apply_alpha_to_image,
    _dependency_hint,
)
from leo_ppt_generator.upstream_bridge import run_upstream
from PIL import Image, ImageDraw
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[2]
EDITABLE_RUNTIME = (
    REPO_ROOT
    / "skills/leo-ppt-generator/runtime/src/leo_ppt_generator/_vendor/editable_ppt/editppt/runtime"
)
if str(EDITABLE_RUNTIME) not in sys.path:
    sys.path.insert(0, str(EDITABLE_RUNTIME))

import _input_normalization as input_normalization
from _input_normalization import normalize_inputs
from formula_renderer import formula_image_fragment
from paddle_text_hints import build_page_hints
from text_hints import page_text_hints
from validate_pptx import quality_contract_violations


def test_confirmed_source_formula_inventory_cannot_be_omitted():
    manifest = {
        "visual_inventory": [],
        "background_strategy": "rebuild",
        "quality_checks": {
            "font_size_calibrated": True,
            "visual_inventory_matched": True,
            "background_strategy_checked": True,
            "shape_corner_geometry_checked": True,
        },
        "expected_formula_inventory": [{"id": "formula_1", "status": "confirmed"}],
        "formula_inventory": [],
    }
    violations = quality_contract_violations(manifest)
    assert any(item.get("formula_id") == "formula_1" for item in violations)
    manifest["formula_inventory"] = [{"id": "formula_1", "decision": "latex-rendered-image"}]
    assert not any(item.get("formula_id") == "formula_1" for item in quality_contract_violations(manifest))


class _OpenAIImages:
    def __init__(self) -> None:
        self.generated = None
        self.edited = None

    def generate(self, **payload):
        self.generated = payload
        return SimpleNamespace(data=[SimpleNamespace(b64_json="Z2VuZXJhdGVk")])

    def edit(self, **payload):
        self.edited = payload
        assert isinstance(payload["image"], list)
        assert len(payload["image"]) == 2
        assert all(not item.closed for item in payload["image"])
        return SimpleNamespace(data=[SimpleNamespace(b64_json="ZWRpdGVk")])


class _AsyncOpenAIImages:
    def __init__(self) -> None:
        self.calls = []

    async def generate(self, **payload):
        self.calls.append(payload)
        return SimpleNamespace(data=[SimpleNamespace(b64_json="YmF0Y2g=")])


def test_openai_compatible_provider_generates_with_forwarded_payload():
    images = _OpenAIImages()
    provider = OpenAICompatibleImageProvider(
        api_key="reference-only", base_url="https://example.test/v1", client_factory=lambda: SimpleNamespace(images=images)
    )
    result = provider.generate({"model": "gpt-image-2", "prompt": "diagram", "size": "1024x1024"})
    assert result == ["Z2VuZXJhdGVk"]
    assert images.generated == {"model": "gpt-image-2", "prompt": "diagram", "size": "1024x1024"}


def test_openai_compatible_provider_edits_with_multiple_reference_images(tmp_path: Path):
    paths = [tmp_path / "target.png", tmp_path / "style.png"]
    for index, path in enumerate(paths):
        Image.new("RGB", (8, 8), (index * 100, 20, 30)).save(path)
    images = _OpenAIImages()
    provider = OpenAICompatibleImageProvider(
        api_key="reference-only", base_url=None, client_factory=lambda: SimpleNamespace(images=images)
    )
    result = provider.edit({"model": "gpt-image-2", "prompt": "preserve target"}, paths, None)
    assert result == ["ZWRpdGVk"]
    assert images.edited["prompt"] == "preserve target"
    assert all(item.closed for item in images.edited["image"])


def test_openai_compatible_provider_generate_batch_uses_async_client():
    images = _AsyncOpenAIImages()
    provider = OpenAICompatibleImageProvider(
        api_key="reference-only",
        base_url=None,
        async_client_factory=lambda: SimpleNamespace(images=images),
    )
    result = asyncio.run(
        provider.generate_batch({"model": "gpt-image-2", "prompt": "one"}, attempts=2, job_label="slide-1")
    )
    assert result == ["YmF0Y2g="]
    assert images.calls == [{"model": "gpt-image-2", "prompt": "one"}]


class _Response:
    def __init__(self, payload: dict | bytes) -> None:
        self.payload = payload

    def __enter__(self):
        return self

    def __exit__(self, *_args):
        return False

    def read(self) -> bytes:
        if isinstance(self.payload, bytes):
            return self.payload
        return json.dumps(self.payload).encode("utf-8")


def test_atlascloud_provider_generates_and_edits_with_data_url_references(tmp_path: Path):
    requests: list[Request] = []
    results = iter(
        [
            {"data": {"id": "generate-1"}},
            {"data": {"status": "completed", "outputs": ["Z2Vu"]}},
            {"data": {"id": "edit-1"}},
            {"data": {"status": "completed", "outputs": ["ZWRpdA=="]}},
        ]
    )

    def urlopen(request: Request, timeout: int):
        assert timeout == 60
        requests.append(request)
        return _Response(next(results))

    source = tmp_path / "source.png"
    Image.new("RGB", (8, 8), "red").save(source)
    provider = AtlasCloudImageProvider(
        api_key="reference-only",
        base_url="https://api.atlascloud.ai/api/v1/model",
        urlopen=urlopen,
        sleep=lambda _seconds: None,
    )
    assert provider.generate({"model": "gpt-image-2", "prompt": "new"}) == ["Z2Vu"]
    assert provider.edit({"model": "gpt-image-2", "prompt": "edit"}, [source], None) == ["ZWRpdA=="]
    generate_payload = json.loads(requests[0].data)
    edit_payload = json.loads(requests[2].data)
    assert generate_payload["model"].endswith("/text-to-image")
    assert edit_payload["model"].endswith("/edit")
    assert edit_payload["images"][0].startswith("data:image/png;base64,")


def test_chroma_key_removal_makes_key_background_transparent():
    image = Image.new("RGBA", (20, 20), (0, 255, 0, 255))
    ImageDraw.Draw(image).rectangle((6, 6, 13, 13), fill=(220, 20, 20, 255))
    transparent = _apply_alpha_to_image(
        image,
        key=(0, 255, 0),
        tolerance=16,
        spill_cleanup=True,
        soft_matte=False,
        transparent_threshold=12,
        opaque_threshold=220,
    )
    assert transparent == 336
    assert image.getpixel((0, 0))[3] == 0
    assert image.getpixel((10, 10))[3] == 255


def test_chroma_key_dependency_hint_is_available():
    hint = _dependency_hint("pillow")
    assert "leo-ppt-generator managed runtime" in hint
    assert str(Path(sys.executable)) in hint
    assert "pip install pillow" in hint
    assert "codex_ppt_runtime.py" not in hint


def test_legacy_ppt_normalization_forwards_requested_dpi(tmp_path: Path, monkeypatch):
    source = tmp_path / "legacy.ppt"
    source.write_bytes(b"legacy fixture")
    observed = {}

    def render_pdf_pages(_pdf, pages_dir, dpi):
        observed["dpi"] = dpi
        page = pages_dir / "page_001/source.png"
        page.parent.mkdir(parents=True)
        Image.new("RGB", (16, 9), "white").save(page)
        return [page]

    def convert_ppt_to_pptx(_source, tmp):
        converted = tmp / "converted.pptx"
        converted.write_bytes(b"converted fixture")
        return converted

    monkeypatch.setattr(input_normalization, "convert_ppt_to_pptx", convert_ppt_to_pptx)
    monkeypatch.setattr(input_normalization, "collect_notes_from_pptx", lambda *_args: [])
    monkeypatch.setattr(input_normalization, "convert_office_to_pdf", lambda _source, tmp: tmp / "converted.pdf")
    monkeypatch.setattr(input_normalization, "render_pdf_pages", render_pdf_pages)

    manifest_path = normalize_inputs([source], job_dir=tmp_path / "run", dpi=144)
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert observed["dpi"] == 144
    assert manifest["input_type"] == "ppt"


def test_image_deck_assembly_writes_speaker_notes(tmp_path: Path):
    images = []
    for number in (1, 2):
        path = tmp_path / f"slide-{number}.png"
        Image.new("RGB", (320, 180), (number * 40, 60, 80)).save(path)
        images.append(str(path))
    out = tmp_path / "deck.pptx"
    assert create_presentation(images, str(out), speaker_notes={1: "first note", 2: "second note"})
    presentation = Presentation(out)
    assert [slide.notes_slide.notes_text_frame.text for slide in presentation.slides] == [
        "first note",
        "second note",
    ]


def test_image_deck_assembly_preserves_source_ratio_by_default(tmp_path: Path):
    source = tmp_path / "wide.png"
    Image.new("RGB", (800, 200), "red").save(source)
    out = tmp_path / "wide-deck.pptx"

    assert create_presentation([str(source)], str(out))
    presentation = Presentation(out)
    picture = next(shape for shape in presentation.slides[0].shapes if shape.shape_type == 13)
    assert picture.width == presentation.slide_width
    assert picture.height < presentation.slide_height


def test_pdf_normalization_rasterizes_every_page_in_order(tmp_path: Path):
    pdf = tmp_path / "source.pdf"
    document = fitz.open()
    for label in ("page-one", "page-two"):
        page = document.new_page(width=320, height=180)
        page.insert_text((30, 50), label)
    document.save(pdf)
    document.close()
    manifest_path = normalize_inputs([pdf], job_dir=tmp_path / "run", dpi=72)
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert manifest["input_type"] == "pdf"
    assert [page["source_page"] for page in manifest["pages"]] == [1, 2]
    assert all((tmp_path / "run" / page["source_image"]).is_file() for page in manifest["pages"])


def test_image_based_pptx_normalization_preserves_page_order_and_notes(tmp_path: Path):
    images = []
    for number in (1, 2):
        path = tmp_path / f"slide-{number}.png"
        Image.new("RGB", (320, 180), (number * 60, 50, 90)).save(path)
        images.append(str(path))
    source = tmp_path / "source.pptx"
    assert create_presentation(images, str(source), speaker_notes={1: "alpha", 2: "beta"})
    manifest_path = normalize_inputs([source], job_dir=tmp_path / "run", dpi=72)
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    notes = json.loads((tmp_path / "run" / manifest["notes_manifest"]).read_text(encoding="utf-8"))
    assert manifest["input_type"] == "pptx"
    assert [page["source_page"] for page in manifest["pages"]] == [1, 2]
    assert [entry["text"] for entry in notes["notes"]] == ["alpha", "beta"]
    assert all(len(entry["text_sha256"]) == 64 for entry in notes["notes"])


def test_builtin_ink_hints_measure_local_text_without_network(tmp_path: Path):
    page_dir = tmp_path / "page"
    page_dir.mkdir()
    image = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(image)
    glyph_height = 30
    glyph_width = 21
    gap = 5
    stroke = 5
    for offset in range(8):
        x = 100 + offset * (glyph_width + gap)
        draw.rectangle(
            (x, 100, x + glyph_width - 1, 100 + glyph_height - 1),
            outline="#1a1a1a",
            width=stroke,
        )
    image.save(page_dir / "source.png")
    hints = page_text_hints(page_dir)
    assert hints["backend"] == "builtin-ink"
    assert hints["lines"]
    assert all("font_pt_if_cjk" in line and "font_pt_if_latin" in line for line in hints["lines"])


def test_paddle_hints_rescale_ocr_boxes_and_preserve_text(tmp_path: Path):
    page_dir = tmp_path / "page"
    page_dir.mkdir()
    image = Image.new("RGB", (400, 200), "white")
    ImageDraw.Draw(image).rectangle((100, 50, 260, 90), fill="black")
    image.save(page_dir / "source.png")
    pruned = {
        "width": 200,
        "height": 100,
        "parsing_res_list": [
            {"block_label": "text", "block_content": "关键结论", "block_bbox": [50, 25, 130, 45]}
        ],
    }
    hints = build_page_hints(page_dir, pruned)
    assert hints["backend"] == "paddleocr-vl"
    assert hints["lines"][0]["text"] == "关键结论"
    assert hints["lines"][0]["box_px"][0] == 100


def test_image_import_records_hash_backend_and_provenance(tmp_path: Path):
    page = tmp_path / "page"
    page.mkdir()
    source = tmp_path / "generated.png"
    Image.new("RGB", (40, 40), "blue").save(source)
    result = run_upstream(
        "editable-ppt",
        [
            "image",
            "import",
            str(page),
            "--job-id",
            "asset-1",
            "--source-image",
            str(source),
            "--dest",
            "assets/imported.png",
            "--role",
            "asset",
            "--backend",
            "openai-compatible-api",
        ],
    )
    assert result["returncode"] == 0, result
    job = json.loads((page / "imagegen-jobs.json").read_text(encoding="utf-8"))["jobs"][0]
    assert job["status"] == "recorded"
    assert job["backend"] == "openai-compatible-api"
    assert len(job["output_sha256"]) == 64


def test_asset_sheet_processing_removes_key_and_splits_assets(tmp_path: Path):
    page = tmp_path / "page"
    assets = page / "assets"
    assets.mkdir(parents=True)
    source = assets / "sheet.png"
    image = Image.new("RGB", (160, 100), "#ff00ff")
    draw = ImageDraw.Draw(image)
    draw.rectangle((15, 20, 55, 80), fill="black")
    draw.rectangle((100, 25, 145, 75), fill="blue")
    image.save(source)
    result = run_upstream(
        "editable-ppt",
        [
            "image",
            "process-sheet",
            str(page),
            "--job-id",
            "icons",
            "--asset-sheet-source",
            "assets/sheet.png",
            "--asset-names",
            "black-icon,blue-icon",
            "--split-min-area",
            "100",
        ],
    )
    assert result["returncode"] == 0, result
    report = json.loads((assets / "icons.split-report.json").read_text(encoding="utf-8"))
    assert [Path(item["path"]).name for item in report["assets"]] == [
        "black-icon.png",
        "blue-icon.png",
    ]
    assert (assets / "icons.asset-sheet-alpha.png").is_file()


def test_page_contact_sheet_compares_source_and_preview(tmp_path: Path):
    page = tmp_path / "page"
    page.mkdir()
    Image.new("RGB", (120, 60), "white").save(page / "source.png")
    Image.new("RGB", (120, 60), "black").save(page / "preview.png")
    result = run_upstream("editable-ppt", ["page", "contact-sheet", str(page)])
    assert result["returncode"] == 0, result
    with Image.open(page / "split_assets_contact.png") as contact:
        assert contact.width > 240
        assert contact.height > 60


def test_formula_fragment_records_rendered_asset_provenance():
    fragment = formula_image_fragment(
        formula_id="f1",
        image_path="assets/f1.svg",
        tex_source="assets/f1.tex",
        box_px="10,20,300,80",
    )
    assert fragment["images"][0]["box_px"] == [10.0, 20.0, 300.0, 80.0]
    assert fragment["asset_provenance"][0]["source_type"] == "latex-rendered-formula"
    assert fragment["formula_inventory"][0]["editable"] is False


def test_editable_image_generate_and_edit_dry_runs_preserve_provider_inputs(tmp_path: Path):
    source = tmp_path / "source.png"
    Image.new("RGB", (40, 40), "white").save(source)
    generated = run_upstream(
        "editable-ppt",
        ["image", "generate", "--prompt", "new icon", "--out", str(tmp_path / "new.png"), "--dry-run"],
    )
    edited = run_upstream(
        "editable-ppt",
        [
            "image",
            "edit",
            "--image",
            str(source),
            "--prompt",
            "preserve shape",
            "--out",
            str(tmp_path / "edited.png"),
            "--dry-run",
        ],
    )
    assert generated["returncode"] == 0, generated
    assert edited["returncode"] == 0, edited
    assert generated["stdout"]["endpoint"] == "/v1/images/generations"
    assert generated["stdout"]["prompt"] == "new icon"
    assert edited["stdout"]["endpoint"] == "/v1/images/edits"
    assert edited["stdout"]["image"] == [str(source)]
    assert edited["stdout"]["prompt"] == "preserve shape"


def test_style_contract_requires_visual_confirmation_and_user_style_precedence():
    contract = (REPO_ROOT / "skills/leo-ppt-generator/references/style-library.md").read_text(encoding="utf-8")
    assert "确认视觉方向前" in contract
    assert "${LEO_PPT_HOME}/styles/<style-name>.md" in contract
    assert "优先读取同名用户" in contract
    assert "重名时先确认覆盖、合并或改名" in contract


def test_page_decision_contract_distinguishes_native_background_and_foreground():
    contract = (REPO_ROOT / "skills/leo-ppt-generator/references/page-decision-tree.md").read_text(encoding="utf-8")
    assert "native" in contract
    assert "background" in contract
    assert "foreground" in contract
    assert "never screenshotted wholesale" in contract


def test_delivery_contract_reports_artifacts_validation_and_not_run_limits():
    skill = (REPO_ROOT / "skills/leo-ppt-generator/SKILL.md").read_text(encoding="utf-8")
    assert "PPTX 与必要逐页/notes/failure report 路径" in skill
    assert "结构验证结果" in skill
    assert "provider/OCR/viewer/desktop/人工视觉验证" in skill
    assert "未运行" in skill
