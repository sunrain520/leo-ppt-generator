from __future__ import annotations

from pathlib import Path

import pytest
from leo_ppt_generator.config.backend_contract import BackendRegistry
from leo_ppt_generator.editable.adapter import EditableAdapter
from leo_ppt_generator.hybrid.assembler import HybridAssembler
from leo_ppt_generator.image_deck.adapter import ImageDeckAdapter
from PIL import Image
from pptx import Presentation

from tests.ppt_fixtures import build_editable_page


def make_sources(root: Path, count: int) -> list[Path]:
    result = []
    for number in range(1, count + 1):
        path = root / f"source-{number}.png"
        Image.new("RGB", (320, 180), (number * 40, 20, 100)).save(path)
        result.append(path)
    return result


def make_editable_artifacts(root: Path, sources: list[Path]):
    adapter = EditableAdapter(root / "editable")
    adapter.prepare(sources, worker_available=True)
    artifacts = []
    for number, source in enumerate(sources, 1):
        page, manifest = build_editable_page(source, root / f"page-{number}.pptx", text=f"editable {number}")
        validation = root / f"validation-{number}.json"
        validation.write_text('{"passed": true}\n', encoding="utf-8")
        artifacts.append(adapter.record(f"page_{number:03d}", page, validation, manifest, expected_revision=number - 1, operation_id=f"editable-{number}", notes=f"note {number}"))
    return artifacts


@pytest.mark.parametrize("backend", ["fixture", "openai"])
def test_generate_route_offline_fixture(tmp_path, backend):
    BackendRegistry.default().select(backend, required={"generate"})
    sources = make_sources(tmp_path, 2)
    adapter = ImageDeckAdapter(tmp_path / "image-run")
    adapter.prepare([{"number": 1, "notes": "note 1"}, {"number": 2, "notes": "note 2"}])
    for number, source in enumerate(sources, 1):
        adapter.record(number, source, backend=backend, expected_revision=number - 1, operation_id=f"image-{number}")
    result = adapter.finalize(tmp_path / f"generate-{backend}.pptx")
    assert result["delivery_type"] == "image"
    assert len(Presentation(result["pptx"]).slides) == 2


@pytest.mark.parametrize("backend", ["fixture", "openai"])
def test_direct_editable_route_offline_fixture(tmp_path, backend):
    BackendRegistry.default().select(backend, required={"edit"})
    artifacts = make_editable_artifacts(tmp_path, make_sources(tmp_path, 2))
    result = HybridAssembler().assemble(artifacts, tmp_path / f"editable-{backend}.pptx")
    assert result["delivery_type"] == "editable"
    assert [page["mode"] for page in result["pages"]] == ["editable", "editable"]


@pytest.mark.parametrize("backend", ["fixture", "openai"])
def test_upgrade_full_preserves_image_delivery_if_editable_stage_fails(tmp_path, backend):
    sources = make_sources(tmp_path, 2)
    image = ImageDeckAdapter(tmp_path / "image-run")
    image.prepare([{"number": 1}, {"number": 2}])
    for number, source in enumerate(sources, 1):
        image.record(number, source, backend=backend, expected_revision=number - 1, operation_id=f"image-{number}")
    original = image.finalize(tmp_path / f"original-{backend}.pptx")
    editable = EditableAdapter(tmp_path / "editable-run")
    editable.prepare(sources, worker_available=True)
    failed_validation = tmp_path / "failed.json"
    failed_validation.write_text('{"passed": false}\n', encoding="utf-8")
    page, manifest = build_editable_page(sources[0], tmp_path / "bad-page.pptx")
    with pytest.raises(Exception, match="page_validation_failed"):
        editable.record("page_001", page, failed_validation, manifest, expected_revision=0, operation_id="failed-editable")
    assert Path(original["pptx"]).is_file()
    assert len(Presentation(original["pptx"]).slides) == 2


@pytest.mark.parametrize("backend", ["fixture", "openai"])
def test_upgrade_selected_and_partial_hybrid_offline_fixture(tmp_path, backend):
    sources = make_sources(tmp_path, 3)
    editable = make_editable_artifacts(tmp_path, [sources[0], sources[2]])
    from leo_ppt_generator.contracts import PageArtifact

    third = PageArtifact.from_source(
        "page_003",
        "editable",
        editable[1].source_path,
        editable[1].artifact_path,
        editable[1].validation_path,
        manifest=editable[1].manifest_path,
        notes=editable[1].notes,
    )
    artifacts = [
        editable[0],
        PageArtifact.from_source("page_002", "image", sources[1], sources[1], None, notes="note 2"),
        third,
    ]
    result = HybridAssembler().assemble(artifacts, tmp_path / f"hybrid-{backend}.pptx", selected_pages={1, 3})
    assert result["delivery_type"] == "hybrid"
    failures = {2: "page_validation_failed"}
    confirmation = HybridAssembler.failure_fingerprint(
        failures,
        selected_pages={1, 2, 3},
        baseline_fingerprint=HybridAssembler.baseline_fingerprint(artifacts),
    )
    partial = HybridAssembler().assemble(
        artifacts,
        tmp_path / f"partial-{backend}.pptx",
        selected_pages={1, 2, 3},
        failures=failures,
        partial_confirmation=confirmation,
    )
    assert partial["delivery_type"] == "partial-hybrid"
    assert partial["pages"][1]["mode"] == "image"
    assert Path(partial["pptx"]).is_file()
